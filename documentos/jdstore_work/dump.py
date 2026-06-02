# -*- coding: utf-8 -*-
import sys
sys.stdout.reconfigure(encoding='utf-8', errors='replace')
from pptx import Presentation
from pptx.util import Emu

p = Presentation(r"C:\Users\berke\Downloads\JD STORE.pptx")
print("DIM:", round(p.slide_width/914400,2), "x", round(p.slide_height/914400,2), "in")
for i, s in enumerate(p.slides, 1):
    texts = []
    imgs = 0
    for sh in s.shapes:
        if sh.shape_type == 13:  # picture
            imgs += 1
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                texts.append(t.replace("\n", " | "))
    print(f"\n## Slide {i}  (imgs={imgs})")
    for t in texts[:8]:
        print("  -", (t[:120] + ("..." if len(t) > 120 else "")))
